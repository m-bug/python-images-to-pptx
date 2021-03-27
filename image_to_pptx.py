##
##  Copy & Paste Tool for images to PowerPoint(.pptx)
##
import pptx
import pptx.util
import glob
import imageio as im

from os import listdir
from PIL import Image as PImage


def loadImages(path):
    # return array of images

    imagesList = listdir(path)
    loadedImages = []
    for image in imagesList:
        img = PImage.open(path + image)
        loadedImages.append(img)

    return loadedImages

path = "./friends/"

# your images in an array
imgs = loadImages(path)

OUTPUT_TAG = "MY_TAG"

# new
prs = pptx.Presentation()
# open
# prs_exists = pptx.Presentation("some_presentation.pptx")

# default slide width
#prs.slide_width = 9144000
# slide height @ 4:3
#prs.slide_height = 6858000
# slide height @ 16:9
prs.slide_height = 5143500

# title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
# blank slide
#slide = prs.slides.add_slide(prs.slide_layouts[6])
layout8 = prs.slide_layouts[8]

# set title
title = slide.shapes.title
title.text = OUTPUT_TAG

pic_left  = int(prs.slide_width * 0.15)
pic_top   = int(prs.slide_height * 0.1)
pic_width = int(prs.slide_width * 0.7)


def _add_image(slide, placeholder_id, image_url):
    placeholder = slide.placeholders[placeholder_id]

    # Calculate the image size of the image
    im = PImage.open(image_url)
    width, height = im.size

    # Make sure the placeholder doesn't zoom in
    placeholder.height = height
    placeholder.width = width

    # Insert the picture
    placeholder = placeholder.insert_picture(image_url)

    # Calculate ratios and compare
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio

    # Placeholder width too wide:
    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side
    # Placeholder height too high
    else:
        difference_on_each_side = -ratio_difference / 2
        placeholder.crop_bottom = -difference_on_each_side
        placeholder.crop_top = -difference_on_each_side

for img in imgs:
    # you can show every image
    #img.show()
    slide = prs.slides.add_slide(layout8)

    title = slide.shapes.title.text = "This is Powerpoint"
    sub = slide.placeholders[2].text = "Python has the power"
    _add_image(slide, 1, img)

prs.save("%s.pptx" % OUTPUT_TAG)